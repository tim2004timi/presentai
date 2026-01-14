import base64
import mimetypes
import aiofiles
from pathlib import Path
from typing import List, Dict, Any
from io import BytesIO
import uuid
import subprocess
import sys

import bcrypt
import pypdf
from fastapi import HTTPException
from openai import AsyncOpenAI
from pptx import Presentation

from src.config import settings
from src.schemas.inputform import InputFormCreate, GeneratedSlides, SlideCard

STATIC_DIR = Path("static")
MODEL_NAME = "gpt-4o-2024-08-06"

async def _encode_image(file_path: Path) -> str:
    async with aiofiles.open(file_path, "rb") as image_file:
        return base64.b64encode(await image_file.read()).decode('utf-8')

async def _read_file_content(file_path: Path) -> str:
    """Умное чтение файла: текст или PDF."""
    try:
        async with aiofiles.open(file_path, "r", encoding="utf-8") as f:
            return await f.read()
    except UnicodeDecodeError:
        if file_path.suffix.lower() == '.pdf':
            try:
                text = ""
                async with aiofiles.open(file_path, "rb") as f:
                    content = await f.read()
                    pdf_reader = pypdf.PdfReader(BytesIO(content))
                    for page in pdf_reader.pages:
                        text += page.extract_text() + "\n"
                return text
            except Exception as e:
                return f"[Ошибка чтения PDF {file_path.name}: {e}]"
        
        return f"[Не удалось прочитать файл {file_path.name}: бинарный формат]"

def _is_image_file(file_path: Path) -> bool:
    mime_type, _ = mimetypes.guess_type(str(file_path))
    if mime_type:
        return mime_type.startswith("image/")
    return file_path.suffix.lower() in {".jpg", ".jpeg", ".png", ".webp"}

async def generate_slides_from_input_form(form: InputFormCreate) -> GeneratedSlides:
    client = AsyncOpenAI(api_key=settings.OPENAI_API_KEY)
    
    content_parts: List[Dict[str, Any]] = [
        {
            "type": "text", 
            "text": f"Название презентации: {form.title}\n\nТЗ пользователя:\n{form.text}"
        }
    ]

    if form.files:
        for filename in form.files:
            file_path = STATIC_DIR / filename
            if not file_path.exists():
                continue

            if _is_image_file(file_path):
                base64_image = await _encode_image(file_path)
                mime_type, _ = mimetypes.guess_type(str(file_path))
                content_parts.append({
                    "type": "image_url",
                    "image_url": {
                        "url": f"data:{mime_type or 'image/jpeg'};base64,{base64_image}"
                    }
                })
            else:
                text_content = await _read_file_content(file_path)
                content_parts.append({
                    "type": "text",
                    "text": f"\n--- Файл {filename} ---\n{text_content}\n--- Конец файла ---\n"
                })

    system_prompt = (
        f"You are a presentation generator. Create exactly {form.slides} slides based on the user content. "
        "Analyze text and provided images/files carefully. Return result in Russian language."
    )

    try:
        completion = await client.beta.chat.completions.parse(
            model=MODEL_NAME,
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": content_parts},
            ],
            response_format=GeneratedSlides, 
            temperature=0.4,
        )
        
        return completion.choices[0].message.parsed

    except Exception as e:
        print(f"OpenAI Error: {e}")
        raise HTTPException(status_code=500, detail="Ошибка генерации контента через AI")


def get_password_hash(password: str) -> str:
    """Генерирует хеш пароля"""
    if isinstance(password, str):
        password = password.encode("utf-8")
    salt = bcrypt.gensalt()
    hashed = bcrypt.hashpw(password, salt)
    return hashed.decode("utf-8")


def verify_password(plain_password: str, hashed_password: str) -> bool:
    """Проверяет соответствие пароля и хеша"""
    if isinstance(hashed_password, str):
        hashed_password = hashed_password.encode("utf-8")
    if isinstance(plain_password, str):
        plain_password = plain_password.encode("utf-8")
    return bcrypt.checkpw(plain_password, hashed_password)


def _set_placeholder_text(slide, placeholder_name: str, text: str) -> None:
    """
    Устанавливает текст в плейсхолдер с указанным именем.
    Использует text_frame для правильной установки текста.
    Для новых слайдов из layout ищет по типу плейсхолдера.
    """
    from pptx.enum.shapes import PP_PLACEHOLDER_TYPE
    
    for ph in slide.placeholders:
        ph_name = getattr(ph, "name", "")
        if ph_name == placeholder_name:
            _set_text_in_placeholder(ph, text)
            return
    
    target_type = None
    if placeholder_name == "Заголовок 1":
        target_type = PP_PLACEHOLDER_TYPE.TITLE
    elif placeholder_name == "Текст 2":
        target_type = PP_PLACEHOLDER_TYPE.BODY
    elif placeholder_name == "Подзаголовок 2":
        target_type = PP_PLACEHOLDER_TYPE.SUBTITLE
    
    if target_type:
        for ph in slide.placeholders:
            if hasattr(ph, "placeholder_format"):
                ph_type = getattr(ph.placeholder_format, "type", None)
                if ph_type == target_type:
                    _set_text_in_placeholder(ph, text)
                    return

    for shape in slide.shapes:
        shape_name = getattr(shape, "name", "")
        if shape_name == placeholder_name:
            _set_text_in_shape(shape, text)
            return


def _set_text_in_placeholder(ph, text: str) -> None:
    """Вспомогательная функция для установки текста в плейсхолдер"""
    if hasattr(ph, "text_frame") and ph.text_frame:
        ph.text_frame.clear()
        ph.text_frame.text = text
        if ph.text_frame.paragraphs:
            ph.text_frame.paragraphs[0].text = text
    elif hasattr(ph, "text"):
        ph.text = text


def _set_text_in_shape(shape, text: str) -> None:
    """Вспомогательная функция для установки текста в shape"""
    if hasattr(shape, "text_frame") and shape.text_frame:
        shape.text_frame.clear()
        shape.text_frame.text = text
        if shape.text_frame.paragraphs:
            shape.text_frame.paragraphs[0].text = text
    elif hasattr(shape, "text"):
        shape.text = text


def generate_pdf(pptx_path: Path) -> Path:
    """
    Конвертирует PPTX файл в PDF используя LibreOffice.
    
    Args:
        pptx_path: Путь к PPTX файлу
        
    Returns:
        Path к созданному PDF файлу
        
    Raises:
        HTTPException: Если конвертация не удалась
    """
    pdf_path = pptx_path.with_suffix('.pdf')
    
    try:
        if sys.platform == "win32":
            libreoffice_paths = [
                r"C:\Program Files\LibreOffice\program\soffice.exe",
                r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
            ]
            soffice = None
            for path in libreoffice_paths:
                if Path(path).exists():
                    soffice = path
                    break
            
            if not soffice:
                raise HTTPException(
                    status_code=500,
                    detail="LibreOffice не найден. Установите LibreOffice для конвертации в PDF.",
                )
            
            cmd = [
                soffice,
                "--headless",
                "--convert-to", "pdf",
                "--outdir", str(pdf_path.parent),
                str(pptx_path),
            ]
        else:
            cmd = [
                "soffice",
                "--headless",
                "--convert-to", "pdf",
                "--outdir", str(pdf_path.parent),
                str(pptx_path),
            ]
        
        result = subprocess.run(
            cmd,
            capture_output=True,
            text=True,
            timeout=60,  # Таймаут 60 секунд
        )
        
        if result.returncode != 0:
            raise HTTPException(
                status_code=500,
                detail=f"Ошибка конвертации PPTX в PDF: {result.stderr}",
            )
        
        if not pdf_path.exists():
            raise HTTPException(
                status_code=500,
                detail="PDF файл не был создан после конвертации",
            )
        
        return pdf_path
        
    except subprocess.TimeoutExpired:
        raise HTTPException(
            status_code=500,
            detail="Таймаут при конвертации PPTX в PDF",
        )
    except FileNotFoundError:
        raise HTTPException(
            status_code=500,
            detail="LibreOffice не найден. Установите LibreOffice для конвертации в PDF.",
        )
    except Exception as e:
        raise HTTPException(
            status_code=500,
            detail=f"Ошибка при конвертации в PDF: {str(e)}",
        )


def generate_pptx(title: str, cards: List[SlideCard]) -> str:
    """
    Генерирует PPTX-презентацию из списка карточек на основе шаблона template.pptx.
    Также создает PDF версию презентации.

    - Слайд 0 шаблона используется как титульный: в плейсхолдер «Заголовок слайда» ставится title.
    - Слайд 1 шаблона используется как основной: для каждой карточки создаётся слайд,
      в «Заголовок слайда» пишется card.title, в «Текст слайда» — card.text.

    Возвращает имя файла (UUID.pptx), сохранённого в директории static.
    """
    template_path = Path("template.pptx")
    if not template_path.exists():
        raise HTTPException(
            status_code=500,
            detail="Не найден файл шаблона презентации template.pptx",
        )

    STATIC_DIR.mkdir(parents=True, exist_ok=True)

    prs = Presentation(str(template_path))

    if len(prs.slides) < 2:
        raise HTTPException(
            status_code=500,
            detail="Шаблон template.pptx должен содержать минимум два слайда (титульный и основной)",
        )

    title_slide = prs.slides[0]
    _set_placeholder_text(title_slide, "Заголовок 1", title)

    base_slide = prs.slides[1]
    base_layout = base_slide.slide_layout

    if len(prs.slides) > 1:
        slide_element = prs.slides._sldIdLst[1]
        slide_rId = slide_element.rId
        
        prs.part.drop_rel(slide_rId)
        
        prs.slides._sldIdLst.remove(slide_element)

    sorted_cards = sorted(cards, key=lambda c: c.index)

    for card in sorted_cards:
        slide = prs.slides.add_slide(base_layout)
        _set_placeholder_text(slide, "Заголовок 1", card.title)
        _set_placeholder_text(slide, "Текст 2", card.text)

    filename = f"{uuid.uuid4()}.pptx"
    output_path = STATIC_DIR / filename
    prs.save(str(output_path))

    try:
        generate_pdf(output_path)
    except HTTPException:
        pass

    return filename

