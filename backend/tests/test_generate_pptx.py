from pathlib import Path

from pptx import Presentation

from src.schemas.inputform import SlideCard
from src.utils import generate_pptx, STATIC_DIR


def test_generate_pptx_creates_file_and_slides():
    # Arrange
    title = "Test"
    cards = [
        SlideCard(index=1, title="Test", text="1. Test\n2. Test\n3. Test"),
    ]

    # Act
    filename = generate_pptx(title, cards)

    # Assert: файл создан в static
    output_path = STATIC_DIR / filename
    assert output_path.exists(), "PPTX файл должен быть создан"
    assert filename.endswith(".pptx")

    # Дополнительно проверим, что структура слайдов соответствует ожиданиям
    prs = Presentation(str(output_path))
    # Ожидаем: титульный + 1 слайд по карточке
    assert len(prs.slides) >= 2

    # Титульный слайд: проверяем, что там есть заголовок
    title_slide = prs.slides[0]
    title_texts = [shape.text for shape in title_slide.shapes if hasattr(shape, "text_frame")]
    assert any("Test" in t for t in title_texts), "Титульный слайд должен содержать заголовок"

    # Слайд с карточкой: проверяем наличие заголовка и текста
    card_slide = prs.slides[1]
    card_texts = [shape.text for shape in card_slide.shapes if hasattr(shape, "text_frame")]
    assert any("Test" in t for t in card_texts), "Слайд карточки должен содержать заголовок"
    assert any("1. Test" in t for t in card_texts), "Слайд карточки должен содержать текст пункта"





